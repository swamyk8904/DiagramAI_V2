from pptx import Presentation


def read_pptx(file_path):
    """
    Reads text from all slides in a PowerPoint presentation.
    """

    try:
        presentation = Presentation(file_path)

        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text.strip()

    except Exception as e:
        return f"Error reading PPTX: {e}"